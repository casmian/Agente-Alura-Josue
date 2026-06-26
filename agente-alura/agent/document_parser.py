import os
import re
import csv
import json
from bs4 import BeautifulSoup
from pypdf import PdfReader
import docx
import openpyxl
from pptx import Presentation

def clean_text_content(text: str) -> str:
    """
    Limpia ruidos, dobles espacios y caracteres corruptos.
    """
    text = text.replace('\r\n', '\n')
    text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F]', '', text)
    text = re.sub(r'[ \t]+', ' ', text)
    text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
    return text.strip()

def parse_html_to_segments(html_content: str) -> list:
    """
    Parsea HTML dividiéndolo por etiquetas H1-H4.
    """
    soup = BeautifulSoup(html_content, 'html.parser')
    for script in soup(["script", "style"]):
        script.decompose()
        
    body = soup.body if soup.body else soup
    segments = []
    current_header = "General"
    current_content = []
    
    for child in body.children:
        if child.name in ['h1', 'h2', 'h3', 'h4']:
            if current_content:
                text = clean_text_content("\n".join(current_content))
                if text:
                    segments.append({"content": text, "location": f"Sección: {current_header}"})
                current_content = []
            current_header = child.get_text().strip()
        elif child.name is not None:
            text = child.get_text().strip()
            if text:
                current_content.append(text)
                
    if current_content:
        text = clean_text_content("\n".join(current_content))
        if text:
            segments.append({"content": text, "location": f"Sección: {current_header}"})
            
    return segments

def parse_pdf_to_segments(file_path: str) -> list:
    """
    Parsea PDF página por página.
    """
    reader = PdfReader(file_path)
    segments = []
    for idx, page in enumerate(reader.pages):
        text = page.extract_text()
        clean = clean_text_content(text or "")
        if clean:
            segments.append({
                "content": clean,
                "location": f"Página {idx + 1}"
            })
    return segments

def parse_docx_to_segments(file_path: str) -> list:
    """
    Parsea Word usando python-docx por encabezados.
    """
    doc = docx.Document(file_path)
    segments = []
    current_header = "General"
    current_content = []
    
    for p in doc.paragraphs:
        if p.style.name.startswith('Heading'):
            if current_content:
                text = clean_text_content("\n".join(current_content))
                if text:
                    segments.append({"content": text, "location": f"Sección: {current_header}"})
                current_content = []
            current_header = p.text.strip()
        else:
            text = p.text.strip()
            if text:
                current_content.append(text)
                
    if current_content:
        text = clean_text_content("\n".join(current_content))
        if text:
            segments.append({"content": text, "location": f"Sección: {current_header}"})
            
    return segments

def parse_xlsx_to_segments(file_path: str, rows_per_segment: int = 10) -> list:
    """
    Parsea Excel fila por fila en grupos.
    """
    wb = openpyxl.load_workbook(file_path, data_only=True)
    segments = []
    
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        rows = list(sheet.rows)
        if not rows:
            continue
            
        headers = [cell.value for cell in rows[0]]
        headers = [str(h) if h is not None else f"Columna_{i}" for i, h in enumerate(headers)]
        
        data_rows = rows[1:]
        for i in range(0, len(data_rows), rows_per_segment):
            row_group = data_rows[i:i + rows_per_segment]
            start_row = i + 2
            end_row = start_row + len(row_group) - 1
            
            row_texts = []
            for idx, row in enumerate(row_group):
                row_num = start_row + idx
                row_details = []
                for col_idx, cell in enumerate(row):
                    header = headers[col_idx] if col_idx < len(headers) else f"Columna_{col_idx}"
                    val = cell.value if cell.value is not None else ""
                    row_details.append(f"{header}: {val}")
                row_texts.append(f"[Fila {row_num}] " + " | ".join(row_details))
                
            content_text = "\n".join(row_texts)
            segments.append({
                "content": content_text,
                "location": f"Hoja: {sheet_name}, Filas {start_row}-{end_row}"
            })
            
    return segments

def parse_pptx_to_segments(file_path: str) -> list:
    """
    Parsea diapositivas de PowerPoint.
    """
    prs = Presentation(file_path)
    segments = []
    for idx, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        
        content = clean_text_content("\n".join(slide_texts))
        if content:
            segments.append({
                "content": content,
                "location": f"Diapositiva {idx + 1}"
            })
    return segments

def parse_csv_to_segments(file_path: str, rows_per_segment: int = 8) -> list:
    """
    Parsea archivos CSV agrupando registros.
    """
    segments = []
    with open(file_path, mode='r', encoding='utf-8') as f:
        reader = csv.DictReader(f)
        rows = list(reader)
        
    for i in range(0, len(rows), rows_per_segment):
        row_group = rows[i:i + rows_per_segment]
        start_row = i + 1
        end_row = start_row + len(row_group) - 1
        
        row_texts = []
        for idx, row in enumerate(row_group):
            row_num = start_row + idx
            details = " | ".join([f"{k}: {v}" for k, v in row.items()])
            row_texts.append(f"[Registro {row_num}] {details}")
            
        content_text = "\n".join(row_texts)
        segments.append({
            "content": content_text,
            "location": f"Filas {start_row}-{end_row}"
        })
        
    return segments

def parse_markdown_to_segments(md_content: str) -> list:
    """
    Parsea Markdown usando encabezados '#'.
    """
    lines = md_content.split('\n')
    segments = []
    current_header = "General"
    current_content = []
    
    for line in lines:
        if line.startswith('#'):
            if current_content:
                text = clean_text_content("\n".join(current_content))
                if text:
                    segments.append({"content": text, "location": f"Sección: {current_header}"})
                current_content = []
            current_header = line.strip()
        else:
            current_content.append(line)
            
    if current_content:
        text = clean_text_content("\n".join(current_content))
        if text:
            segments.append({"content": text, "location": f"Sección: {current_header}"})
            
    return segments

def parse_json_to_segments(json_content: str) -> list:
    """
    Parsea JSON estructurado.
    """
    obj = json.loads(json_content)
    segments = []
    if isinstance(obj, list):
        for index, item in enumerate(obj):
            segments.append({
                "content": clean_text_content(json.dumps(item, indent=2, ensure_ascii=False)),
                "location": f"Registro índice {index}"
            })
    elif isinstance(obj, dict):
        for key, val in obj.items():
            segments.append({
                "content": clean_text_content(f"{key}: {json.dumps(val, indent=2, ensure_ascii=False)}"),
                "location": f"Llave principal: {key}"
            })
    return segments

def parse_document_to_segments(file_path: str) -> list:
    """
    Despacha el parseador correcto según la extensión.
    """
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"El archivo no existe en: {file_path}")
        
    if ext in ['.md', '.txt']:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        return parse_markdown_to_segments(content)
        
    elif ext == '.json':
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        return parse_json_to_segments(content)
        
    elif ext == '.html':
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        return parse_html_to_segments(content)
        
    elif ext == '.pdf':
        return parse_pdf_to_segments(file_path)
        
    elif ext == '.docx':
        return parse_docx_to_segments(file_path)
        
    elif ext in ['.xlsx', '.xls']:
        return parse_xlsx_to_segments(file_path)
        
    elif ext == '.pptx':
        return parse_pptx_to_segments(file_path)
        
    elif ext == '.csv':
        return parse_csv_to_segments(file_path)
        
    else:
        raise ValueError(f"Extensión no soportada: {ext}")
